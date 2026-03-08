#!/usr/bin/env python3
"""
PPTX Generator - Convert YAML files to PowerPoint presentations

This script converts YAML configuration files into PowerPoint presentations.
It supports both single file conversion and batch processing of multiple files.
"""

import yaml
import os
import argparse
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


class PPTXGenerator:
    """Main class for generating PowerPoint presentations from YAML files."""
    
    def __init__(self, default_layout=1, verbose=False):
        """
        Initialize the PPTX generator.
        
        Args:
            default_layout (int): Default slide layout number
            verbose (bool): Enable verbose output
        """
        self.default_layout = default_layout
        self.verbose = verbose
        
    def log(self, message):
        """Print verbose messages if verbose mode is enabled."""
        if self.verbose:
            print(f"[INFO] {message}")
    
    def validate_yaml_config(self, config, yaml_file):
        """Validate the YAML configuration structure."""
        required_fields = ['title', 'slides']
        missing_fields = [field for field in required_fields if field not in config]
        
        if missing_fields:
            raise ValueError(f"Missing required fields in {yaml_file}: {', '.join(missing_fields)}")
        
        if not isinstance(config['slides'], list):
            raise ValueError(f"'slides' must be a list in {yaml_file}")
        
        for i, slide in enumerate(config['slides']):
            if 'title' not in slide or 'content' not in slide:
                raise ValueError(f"Slide {i+1} missing 'title' or 'content' in {yaml_file}")
            
            if not isinstance(slide['content'], list):
                raise ValueError(f"Slide {i+1} 'content' must be a list in {yaml_file}")
    
    def load_yaml_config(self, yaml_file):
        """Load and validate YAML configuration from file."""
        try:
            with open(yaml_file, 'r', encoding='utf-8') as f:
                config = yaml.safe_load(f)
            
            if config is None:
                raise ValueError(f"Empty YAML file: {yaml_file}")
            
            self.validate_yaml_config(config, yaml_file)
            return config
            
        except yaml.YAMLError as e:
            raise ValueError(f"Invalid YAML syntax in {yaml_file}: {e}")
        except FileNotFoundError:
            raise FileNotFoundError(f"YAML file not found: {yaml_file}")
        except Exception as e:
            raise Exception(f"Error loading {yaml_file}: {e}")
    
    def create_presentation(self, config, layout=None):
        """Create a PowerPoint presentation from configuration."""
        prs = Presentation()
        slide_layout = prs.slide_layouts[layout or self.default_layout]
        
        for slide_data in config['slides']:
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            if slide.shapes.title:
                slide.shapes.title.text = slide_data['title']
            
            # Set content
            if len(slide.placeholders) > 1:
                content_text = "\n".join(slide_data['content'])
                slide.placeholders[1].text = content_text
        
        return prs
    
    def generate_single_ppt(self, yaml_file, output_file=None, layout=None):
        """Generate a single PowerPoint presentation from a YAML file."""
        self.log(f"Processing: {yaml_file}")
        
        # Load configuration
        config = self.load_yaml_config(yaml_file)
        
        # Determine output file path
        if output_file is None:
            output_file = config.get('output', f"{Path(yaml_file).stem}.pptx")
        
        # Ensure output directory exists
        output_path = Path(output_file)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        # Create presentation
        prs = self.create_presentation(config, layout)
        
        # Save presentation
        prs.save(str(output_path))
        
        self.log(f"Generated: {output_file}")
        return str(output_path)
    
    def generate_batch_ppt(self, input_dir, output_dir, layout=None):
        """Generate PowerPoint presentations from all YAML files in a directory."""
        input_path = Path(input_dir)
        output_path = Path(output_dir)
        
        if not input_path.exists():
            raise FileNotFoundError(f"Input directory not found: {input_dir}")
        
        output_path.mkdir(parents=True, exist_ok=True)
        
        yaml_files = list(input_path.glob("*.yml")) + list(input_path.glob("*.yaml"))
        
        if not yaml_files:
            self.log(f"No YAML files found in {input_dir}")
            return []
        
        generated_files = []
        
        for yaml_file in yaml_files:
            try:
                output_file = output_path / f"{yaml_file.stem}.pptx"
                result = self.generate_single_ppt(str(yaml_file), str(output_file), layout)
                generated_files.append(result)
            except Exception as e:
                self.log(f"Error processing {yaml_file}: {e}")
                continue
        
        self.log(f"Batch processing complete. Generated {len(generated_files)} files.")
        return generated_files


def parse_arguments():
    """Parse command line arguments."""
    parser = argparse.ArgumentParser(
        description="Convert YAML files to PowerPoint presentations",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Single file conversion
  python generate_ppt.py --input presentations/demo.yml --output output/demo.pptx
  
  # Batch processing
  python generate_ppt.py --input-dir presentations/ --output-dir output/
  
  # Custom layout
  python generate_ppt.py --input presentation.yml --output output.pptx --layout 2
  
  # Verbose output
  python generate_ppt.py --input presentation.yml --output output.pptx --verbose
        """
    )
    
    # Input options
    input_group = parser.add_mutually_exclusive_group(required=False)
    input_group.add_argument(
        '--input', '-i',
        help='Input YAML file path'
    )
    input_group.add_argument(
        '--input-dir', '-d',
        help='Input directory containing YAML files'
    )
    
    # Output options
    output_group = parser.add_mutually_exclusive_group(required=False)
    output_group.add_argument(
        '--output', '-o',
        help='Output PPTX file path (for single file conversion)'
    )
    output_group.add_argument(
        '--output-dir', '-D',
        help='Output directory for PPTX files (for batch processing)'
    )
    
    # Other options
    parser.add_argument(
        '--layout', '-l',
        type=int,
        default=1,
        choices=range(1, 12),
        help='Slide layout number (1-11, default: 1)'
    )
    parser.add_argument(
        '--verbose', '-v',
        action='store_true',
        help='Enable verbose output'
    )
    
    return parser.parse_args()


def main():
    """Main function to handle command line execution."""
    args = parse_arguments()
    
    # Initialize generator
    generator = PPTXGenerator(default_layout=args.layout, verbose=args.verbose)
    
    try:
        # Single file processing
        if args.input:
            if not os.path.isfile(args.input):
                print(f"Error: Input file not found: {args.input}", file=sys.stderr)
                sys.exit(1)
            
            output_file = args.output
            result = generator.generate_single_ppt(args.input, output_file, args.layout)
            print(f"Successfully generated: {result}")
        
        # Batch processing
        elif args.input_dir:
            output_dir = args.output_dir or "output"
            results = generator.generate_batch_ppt(args.input_dir, output_dir, args.layout)
            
            if results:
                print(f"Successfully generated {len(results)} presentations:")
                for result in results:
                    print(f"  - {result}")
            else:
                print("No presentations were generated.", file=sys.stderr)
                sys.exit(1)
        
        # Default behavior (backward compatibility)
        else:
            print("Using default behavior (processing inputs/ directory)...")
            results = generator.generate_batch_ppt("inputs", "outputs", args.layout)
            
            if results:
                print(f"Successfully generated {len(results)} presentations:")
                for result in results:
                    print(f"  - {result}")
            else:
                print("No YAML files found in 'pptx' directory.")
    
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()